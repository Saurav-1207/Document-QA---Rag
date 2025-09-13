import streamlit as st
from pinecone import Pinecone, ServerlessSpec
import cohere
from sentence_transformers import SentenceTransformer
import fitz  # PyMuPDF
import docx  # python-docx
import pandas as pd
from pptx import Presentation

# ---------- Initialize Pinecone ----------
pc = Pinecone(api_key="5b4097f7-41fe-4f3e-8485-b021b82f707a")
index_name = "rag-qa-bot"
if index_name not in pc.list_indexes().names():
    pc.create_index(
        name=index_name,
        dimension=384,
        metric="cosine",
        spec=ServerlessSpec(cloud="aws", region="us-east-1"),
    )
index = pc.Index(index_name)

# ---------- Initialize Cohere Client (v2) ----------
co = cohere.ClientV2(api_key="DSka7LPCNDbFEPvEu6WtAvwZ1iGeoK8OOQi0dsPC")

# ---------- Load Embedding Model ----------
model = SentenceTransformer("all-MiniLM-L12-v2")  # 384 dims

# ---------- File Reading Helpers ----------
@st.cache_data
def read_pdf(file):
    pdf_reader = fitz.open(stream=file.read(), filetype="pdf")
    return " ".join(page.get_text("text") for page in pdf_reader)


@st.cache_data
def read_docx(file):
    doc = docx.Document(file)
    return "\n".join(para.text for para in doc.paragraphs)


@st.cache_data
def read_csv(file):
    df = pd.read_csv(file)
    return df.astype(str).agg(" ".join, axis=1).str.cat(sep=" ")


@st.cache_data
def read_pptx(file):
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


# ---------- Cohere Chat Answer ----------
@st.cache_data
def generate_answer(query, relevant_chunks):
    context = " ".join(relevant_chunks)
    response = co.chat(
        model="command-r-plus",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that answers questions using the provided context."},
            {"role": "user", "content": f"Context: {context}\n\nQuestion: {query}\n\nAnswer:"},
        ],
        temperature=0.7,
    )
    return response.message.content[0].text.strip()


# ---------- Pinecone Query ----------
@st.cache_data
def answer_question(query, namespace):
    query_embedding = model.encode(query).tolist()
    result = index.query(vector=query_embedding, top_k=5, namespace=namespace, include_metadata=True)
    if result["matches"]:
        return [match["metadata"]["content"] for match in result["matches"]]
    return []


# ---------- Delete Namespace ----------
def delete_uploaded_text(namespace):
    index.delete(delete_all=True, namespace=namespace)
    st.session_state["uploaded_files"].remove(namespace)
    st.success(f"Deleted all documents under namespace '{namespace}'.")


# ---------- Streamlit UI ----------
st.set_page_config(page_title="Document QA Chatbot", layout="wide")
st.title("📄 Document QA Chatbot")
# (CSS styling code stays the same...)

# ---- File Upload Sidebar ----
with st.sidebar:
    st.header("Upload Document")
    uploaded_file = st.file_uploader("Choose a PDF, Word, PPTX, or CSV file", type=["pdf", "docx", "pptx", "txt", "csv"])
    st.write("Uploaded Files:")
    if "uploaded_files" not in st.session_state:
        st.session_state["uploaded_files"] = []
    for file in st.session_state["uploaded_files"]:
        st.write(file)

    if uploaded_file is not None and uploaded_file.name not in st.session_state["uploaded_files"]:
        namespace = uploaded_file.name
        if uploaded_file.type == "application/pdf":
            document_text = read_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            document_text = read_docx(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            document_text = read_pptx(uploaded_file)
        elif uploaded_file.type == "text/plain":
            document_text = uploaded_file.read().decode("utf-8")
        elif uploaded_file.type == "text/csv":
            document_text = read_csv(uploaded_file)

        st.session_state["uploaded_files"].append(namespace)
        document_chunks = [document_text[i:i+500] for i in range(0, len(document_text), 500)]
        st.session_state[f"{namespace}_chunks"] = document_chunks

        batch_size = 10
        for i in range(0, len(document_chunks), batch_size):
            batch = document_chunks[i:i+batch_size]
            vectors = [{
                "id": f"chunk-{j+i}",
                "values": model.encode(chunk).tolist(),
                "metadata": {"content": chunk},
            } for j, chunk in enumerate(batch)]
            index.upsert(vectors=vectors, namespace=namespace)

        st.success(f"Document '{uploaded_file.name}' indexed successfully!")


# ---------- Chat State ----------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "relevant_segments" not in st.session_state:
    st.session_state.relevant_segments = []
if "last_processed_input" not in st.session_state:
    st.session_state.last_processed_input = ""


# ---------- Handle Input ----------
def handle_submit():
    user_input = st.session_state.user_input
    if user_input and user_input != st.session_state.last_processed_input:
        st.session_state.last_processed_input = user_input
        st.session_state.chat_history.append({"role": "user", "text": user_input})

        if st.session_state["uploaded_files"]:
            current_namespace = st.session_state["uploaded_files"][-1]
            with st.spinner("Generating answer..."):
                relevant_chunks = answer_question(user_input, current_namespace)
                if relevant_chunks:
                    answer = generate_answer(user_input, relevant_chunks)
                    st.session_state.chat_history.append({"role": "bot", "text": answer})
                    st.session_state.relevant_segments = relevant_chunks
                else:
                    st.session_state.chat_history.append({"role": "bot", "text": "No relevant information found."})
        else:
            st.session_state.chat_history.append({"role": "bot", "text": "No documents uploaded."})


# ---------- Chat UI ----------
st.subheader("Chat with your document:")

def display_chat():
    with st.container():
        st.markdown("<div class='chat-container'>", unsafe_allow_html=True)
        for msg in st.session_state.chat_history:
            role_class = "user-message" if msg["role"] == "user" else "bot-message"
            st.markdown(f"<div class='{role_class}'>{msg['text']}</div>", unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)

display_chat()

with st.form(key="chat_form", clear_on_submit=True):
    st.text_input("You:", key="user_input")
    st.form_submit_button("Send", on_click=handle_submit)

if st.session_state.chat_history:
    with st.sidebar:
        if st.button("Clear Chat"):
            st.session_state.chat_history = []
            st.session_state.relevant_segments = []
            st.session_state.last_processed_input = ""

if st.session_state.relevant_segments:
    st.subheader("Relevant Document Segments:")
    for segment in st.session_state.relevant_segments:
        st.markdown(f"<div class='relevant-segment'>{segment}</div>", unsafe_allow_html=True)
